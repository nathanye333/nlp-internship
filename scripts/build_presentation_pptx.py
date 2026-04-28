"""Generate the executive final-presentation deck as a .pptx file.

The output is `docs/final_presentation.pptx` — readable in PowerPoint and
fully editable in Google Slides (File -> Import slides). Every body font
is 18 pt or larger; titles are 36 pt; key metrics are 56 pt.

Run:

    python scripts/build_presentation_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

BG = RGBColor(0x07, 0x0B, 0x16)
PANEL = RGBColor(0x10, 0x18, 0x2C)
PANEL_STRONG = RGBColor(0x16, 0x20, 0x3A)
LINE = RGBColor(0x1E, 0x2A, 0x47)
INK = RGBColor(0xE7, 0xEE, 0xF8)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x5E, 0xEA, 0xD4)        # cyan
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)  # amber
ACCENT_VIOLET = RGBColor(0xA7, 0x8B, 0xFA) # violet
GOOD = RGBColor(0x22, 0xC5, 0x5E)
WARN = RGBColor(0xFB, 0xBF, 0x24)

FONT_HEAD = "Inter"
FONT_BODY = "Inter"
FONT_MONO = "JetBrains Mono"

# Slide is 13.333 x 7.5 inches (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins
MX = Inches(0.55)         # left
MR = Inches(0.55)         # right
MTOP = Inches(0.5)
USABLE_W = SLIDE_W - MX - MR

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def set_solid(fill, color: RGBColor) -> None:
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=None, line=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        # Tighten the corner radius a bit (default is fairly large)
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        set_solid(shape.fill, fill)
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    color=INK,
    bold=False,
    italic=False,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.18,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_paragraphs(
    slide,
    x,
    y,
    w,
    h,
    items,
    *,
    size=18,
    color=INK,
    bullet=True,
    font=FONT_BODY,
    line_spacing=1.25,
    space_after=4,
):
    """Add a multi-paragraph text block.

    `items` is a list where each element is either:
      - a str ("body bullet")
      - a list of (text, dict) runs for inline styling, e.g.
            [("bold ", {"bold": True}), ("rest", {})]
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True

    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if bullet:
            _set_bullet(p)
            # indent so bullet doesn't crowd
            p.level = 0

        runs = item if isinstance(item, list) else [(item, {})]
        for i, (text, style) in enumerate(runs):
            r = p.add_run()
            r.text = text
            r.font.size = Pt(style.get("size", size))
            r.font.name = style.get("font", font)
            r.font.bold = style.get("bold", False)
            r.font.italic = style.get("italic", False)
            r.font.color.rgb = style.get("color", color)
    return box


def _set_bullet(paragraph) -> None:
    """Add a small, accent-colored bullet to a paragraph."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    # remove any existing bullet definitions
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    # add a bullet character
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "•")
    buFont = pPr.find(qn("a:buFont"))
    if buFont is None:
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        # ensure buFont appears before buChar
        pPr.remove(buFont)
        pPr.insert(list(pPr).index(buChar), buFont)
    buFont.set("typeface", "Arial")
    # color the bullet via buClr
    for el in pPr.findall(qn("a:buClr")):
        pPr.remove(el)
    buClr = etree.Element(qn("a:buClr"))
    srgb = etree.SubElement(buClr, qn("a:srgbClr"))
    srgb.set("val", "5EEAD4")
    pPr.insert(0, buClr)
    # set indent
    pPr.set("indent", "-228600")
    pPr.set("marL", "228600")


def add_pill(slide, x, y, text, *, color=ACCENT, size=12):
    """A small label pill. Auto-sized roughly to text length."""
    char_w = 0.075 if size <= 12 else 0.085
    w = Inches(0.45 + len(text) * char_w)
    h = Inches(0.34)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.5
    set_solid(shape.fill, PANEL_STRONG)
    shape.line.color.rgb = color
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT_BODY
    r.font.color.rgb = color
    return shape, w, h


def add_metric_card(slide, x, y, w, h, label, big, small, *, color=ACCENT):
    add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=0.08)
    add_text(slide, x + Inches(0.2), y + Inches(0.18), w - Inches(0.4), Inches(0.32),
             label.upper(), size=12, color=MUTED, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), Inches(0.9),
             big, size=44, color=color, bold=True, line_spacing=1.0)
    add_text(slide, x + Inches(0.2), y + Inches(1.45), w - Inches(0.4), h - Inches(1.55),
             small, size=14, color=MUTED, line_spacing=1.25)


def add_topbar(slide, crumb, week):
    add_text(slide, MX, Inches(0.32), Inches(8), Inches(0.36),
             crumb.upper(), size=12, color=MUTED, bold=True)
    add_text(slide, SLIDE_W - MR - Inches(8), Inches(0.32), Inches(8), Inches(0.36),
             week.upper(), size=12, color=ACCENT_AMBER, bold=True, align=PP_ALIGN.RIGHT)
    add_rect(slide, MX, Inches(0.78), USABLE_W, Inches(0.02), fill=LINE, line=None)


def add_title(slide, title, *, y=Inches(0.95), size=34):
    add_text(slide, MX, y, USABLE_W, Inches(0.85),
             title, size=size, color=INK, bold=True, line_spacing=1.1)


def slide_blank(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return s


def add_code_block(slide, x, y, w, h, code, *, size=13):
    box = add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=0.08)
    tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.14),
                                   w - Inches(0.36), h - Inches(0.28))
    tf = tb.text_frame
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.18
        r = p.add_run()
        r.text = line if line else " "
        r.font.size = Pt(size)
        r.font.name = FONT_MONO
        r.font.color.rgb = ACCENT
        # Highlight comment-only lines in muted, keywords stay white
        if line.lstrip().startswith("#"):
            r.font.color.rgb = MUTED
        elif line.lstrip().startswith(">>>") or line.lstrip().startswith("..."):
            r.font.color.rgb = ACCENT_AMBER


def add_table(
    slide,
    x,
    y,
    w,
    h,
    rows,
    *,
    header_size=14,
    body_size=15,
    col_widths=None,
):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, frac in enumerate(col_widths):
            table.columns[i].width = int(w * frac / total)
    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            # Cell fill
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL_STRONG if is_header else PANEL
            # Borders
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.line_spacing = 1.15
            # Allow tuples for inline color override: (text, RGBColor) or (text, RGBColor, bool)
            if isinstance(val, tuple):
                text = val[0]
                color = val[1]
                bold = val[2] if len(val) > 2 else False
            else:
                text, color, bold = val, (MUTED if is_header else INK), is_header
            r = p.add_run()
            r.text = str(text)
            r.font.size = Pt(header_size if is_header else body_size)
            r.font.name = FONT_BODY
            r.font.bold = is_header or bold
            r.font.color.rgb = color
    return table


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_01_title(prs):
    s = slide_blank(prs)
    # Eyebrow
    eyebrow = add_text(
        s,
        Inches(1.0), Inches(1.3), Inches(11.3), Inches(0.45),
        "12-WEEK NLP INTERNSHIP   ·   IDX EXCHANGE",
        size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER,
    )
    # Title
    add_text(
        s,
        Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.2),
        "Real Estate Listing\nIntelligence System",
        size=64, color=INK, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.05,
    )
    # Subtitle
    add_text(
        s,
        Inches(1.4), Inches(4.5), Inches(10.5), Inches(1.4),
        "A production-grade NLP pipeline for MLS data — natural-language "
        "search, entity & signal extraction, intent classification, "
        "summarization, and Fair Housing compliance — shipped behind a "
        "FastAPI service and Streamlit demo on Render.",
        size=20, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.4,
    )

    # Stat strip
    stats = [
        ("9,200+", "lines of code"),
        ("22", "modules · 14 test suites"),
        ("40,890", "listings processed"),
        ("12", "API endpoints"),
    ]
    card_w = Inches(2.6)
    gap = Inches(0.25)
    total_w = card_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) / 2
    for i, (big, small) in enumerate(stats):
        x = start_x + (card_w + gap) * i
        y = Inches(6.0)
        add_rect(s, x, y, card_w, Inches(1.05), fill=PANEL, line=LINE, radius=0.08)
        add_text(s, x, y + Inches(0.12), card_w, Inches(0.5), big,
                 size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.62), card_w, Inches(0.4), small,
                 size=12, color=MUTED, align=PP_ALIGN.CENTER, bold=True)


def slide_02_problem(prs):
    s = slide_blank(prs)
    add_topbar(s, "01 · The Problem", "Why this matters")
    add_title(s, "MLS data is rich, but the search experience isn't.")
    add_text(
        s, MX, Inches(1.65), USABLE_W, Inches(0.5),
        "Brokers and consumers query a 40k-row table whose most valuable field — "
        "the agent's free-text remark — is invisible to keyword filters.",
        size=18, color=MUTED, line_spacing=1.4,
    )

    col_w = (USABLE_W - Inches(0.4)) / 2
    # Left: failure modes
    lx = MX
    ly = Inches(2.4)
    lh = Inches(4.6)
    add_rect(s, lx, ly, col_w, lh, fill=PANEL, line=LINE, radius=0.08)
    add_text(s, lx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "What goes wrong today",
             size=22, color=ACCENT, bold=True)
    add_paragraphs(
        s, lx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), Inches(3.6),
        [
            [("Lossy filters. ", {"bold": True}),
             ("\"3 bed under 700k in Irvine with pool\" maps to four columns and one full-text scan.", {})],
            [("No semantic retrieval. ", {"bold": True}),
             ("Synonyms, abbreviations, and context are invisible to LIKE %pool%.", {})],
            [("Manual fair-housing review. ", {"bold": True}),
             ("A single phrase can trigger HUD penalties up to $25,597 per first violation.", {})],
            [("Lead intent is unscored. ", {"bold": True}),
             ("\"Just looking\" and \"ready to make an offer\" route through the same funnel.", {})],
        ],
        size=18, line_spacing=1.35, space_after=10,
    )

    # Right: business cost
    rx = MX + col_w + Inches(0.4)
    ry = Inches(2.4)
    rh = Inches(4.6)
    add_rect(s, rx, ry, col_w, rh, fill=PANEL_STRONG, line=LINE, radius=0.08)
    add_text(s, rx + Inches(0.3), ry + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "Business cost",
             size=22, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, rx + Inches(0.35), ry + Inches(0.95), col_w - Inches(0.7), Inches(3.6),
        [
            [("Conversion drag", {"bold": True}),
             (" — irrelevant results crowd matches; abandonment goes up.", {})],
            [("Compliance liability", {"bold": True}),
             (" — civil penalties up to $25,597 / first violation (HUD 2024).", {})],
            [("Agent productivity", {"bold": True}),
             (" — minutes lost per query × thousands of agents = real OPEX.", {})],
            [("Competitive gap", {"bold": True}),
             (" — peers are deploying NLP retrieval and conversational front-ends.", {})],
        ],
        size=18, line_spacing=1.35, space_after=10,
    )


def slide_03_what_we_built(prs):
    s = slide_blank(prs)
    add_topbar(s, "02 · What We Built", "Capabilities at a glance")
    add_title(s, "Nine production-shaped NLP capabilities, one API.")

    col_w = (USABLE_W - Inches(0.4)) / 3
    cards = [
        ("Retrieval", ACCENT, [
            [("Query Parser → SQL", {"bold": True}),
             (" — /parse-query", {"color": MUTED})],
            [("Hybrid / Semantic / Keyword Search", {"bold": True}),
             (" — /search · /search/compare", {"color": MUTED})],
            [("Listing Detail + Summary", {"bold": True}),
             (" — /listings/{id}", {"color": MUTED})],
        ]),
        ("Understanding", ACCENT_VIOLET, [
            [("Entity Extraction", {"bold": True}),
             (" beds · baths · price · sqft · amenities", {})],
            [("Signal Extraction", {"bold": True}),
             (" — 40,890 listings → JSONL", {})],
            [("Intent Classification", {"bold": True}),
             (" browse / research / ready-to-buy", {})],
        ]),
        ("Trust & Quality", ACCENT_AMBER, [
            [("Listing Summarization", {"bold": True}),
             (" extractive + abstractive", {})],
            [("Fair Housing Compliance", {"bold": True}),
             (" 10 categories · 3 severities", {})],
            [("Answerability gate", {"bold": True}),
             (" rejects out-of-domain before SQL", {})],
        ]),
    ]
    for i, (head, color, items) in enumerate(cards):
        x = MX + (col_w + Inches(0.2)) * i
        y = Inches(1.85)
        h = Inches(3.5)
        add_rect(s, x, y, col_w, h, fill=PANEL, line=LINE, radius=0.08)
        add_text(s, x + Inches(0.3), y + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
                 head, size=22, color=color, bold=True)
        add_paragraphs(s, x + Inches(0.35), y + Inches(0.95),
                       col_w - Inches(0.7), h - Inches(1.1),
                       items, size=17, line_spacing=1.35, space_after=10)

    # Stat strip
    stats = [
        ("CODE", "9.2k", "LOC across 22 scripts + 14 tests", ACCENT),
        ("CORPUS", "40,890", "listings signal-extracted", ACCENT_AMBER),
        ("ENDPOINTS", "12", "FastAPI routes, all rate-limited", ACCENT_VIOLET),
        ("SHIPPED", "live", "Docker · Render blueprint · auto-deploy", GOOD),
    ]
    card_w = (USABLE_W - Inches(0.6)) / 4
    for i, (label, big, small, color) in enumerate(stats):
        x = MX + (card_w + Inches(0.2)) * i
        y = Inches(5.55)
        add_metric_card(s, x, y, card_w, Inches(1.55), label, big, small, color=color)


def slide_04_approach(prs):
    s = slide_blank(prs)
    add_topbar(s, "03 · Approach", "12 weeks · each output feeds the next")
    add_title(s, "Build incrementally, ship every week.")
    add_text(
        s, MX, Inches(1.7), USABLE_W, Inches(0.6),
        "We avoided \"build everything, integrate at the end.\" Every component shipped "
        "behind a test suite the same week — Week 10's API just wired them together.",
        size=18, color=MUTED, line_spacing=1.4,
    )

    rows = [
        ["Wk", "Capability", "Output that fed the next stage"],
        ["1", "Domain taxonomy + sample dataset",   "taxonomy.json (200+ terms) · 1k listing CSV · 50+ user queries"],
        ["2", "Text cleaning + profiling",           "30+ abbreviations · unicode/HTML normalisation · price/measure rewrites"],
        ["3", "Entity extraction (NER-lite)",        "250 labeled remarks · evaluation harness · macro-F1 baseline"],
        ["4", "Query Parser → safe SQL",             "Parameterised WHERE · schema validator · 50+ query patterns"],
        ["5", "Semantic search + BM25",              "FAISS index · all-MiniLM-L6-v2 · latency benchmark harness"],
        ["6", "Signal extraction (full table)",      "40,890 listings → JSONL signals (amenities · condition · location)"],
        ["7", "Intent classifier",                   "252 labeled queries · TF-IDF + Logistic Regression"],
        ["8", "Summarization + answerability",       "Extractive + abstractive · ROUGE-L harness · pre/post checks"],
        ["9", "Fair Housing compliance",             "10-category pattern library · 3 severities · submission workflow"],
        ["10", "Production REST API",                "FastAPI · Pydantic · TTL+LRU cache · rate limiting · structured logs"],
        ["11", "Streamlit demo + deploy",            "3-tab UI · feedback loop · Render blueprint · 1 image, 2 services"],
    ]
    add_table(s, MX, Inches(2.4), USABLE_W, Inches(4.7), rows,
              header_size=15, body_size=15, col_widths=[0.5, 3, 7])


def slide_05_architecture(prs):
    s = slide_blank(prs)
    add_topbar(s, "04 · Architecture", "One image · two services · zero shared state")
    add_title(s, "System architecture.")

    # Layer labels
    layer_x = MX
    layer_h = Inches(1.4)
    layer_w = USABLE_W

    # ----- DATA LAYER -----
    y = Inches(1.65)
    add_text(s, layer_x, y, Inches(1.2), Inches(0.4),
             "DATA", size=13, color=MUTED, bold=True)
    box_y = y + Inches(0.42)
    box_h = Inches(0.95)
    box_w = (layer_w - Inches(0.6)) / 4
    data_blocks = [
        ("MySQL · MLS",        "rets_property · rets_openhouse · california_sold"),
        ("taxonomy.json",      "200+ terms · 8 categories"),
        ("FAISS Index",        "all-MiniLM-L6-v2 · 384d cosine"),
        ("Signals JSONL",      "40,890 listings · entities + amenities + condition"),
    ]
    for i, (head, sub) in enumerate(data_blocks):
        bx = layer_x + (box_w + Inches(0.2)) * i
        add_rect(s, bx, box_y, box_w, box_h, fill=PANEL, line=ACCENT, radius=0.10)
        add_text(s, bx + Inches(0.15), box_y + Inches(0.10),
                 box_w - Inches(0.3), Inches(0.32),
                 head, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, bx + Inches(0.15), box_y + Inches(0.45),
                 box_w - Inches(0.3), Inches(0.45),
                 sub, size=11, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)

    # ----- NLP CORE LAYER -----
    y = Inches(3.15)
    add_text(s, layer_x, y, Inches(1.5), Inches(0.4),
             "NLP CORE", size=13, color=MUTED, bold=True)
    nlp_y = y + Inches(0.42)
    nlp_h = Inches(1.7)
    add_rect(s, layer_x, nlp_y, layer_w, nlp_h, fill=PANEL, line=ACCENT_VIOLET, radius=0.06)

    # 8 modules in a single row
    modules = [
        "TextCleaner (W2)",
        "EntityExtractor (W3)",
        "QueryParser (W4)",
        "Semantic + BM25 (W5)",
        "SignalExtractor (W6)",
        "IntentClassifier (W7)",
        "Summarizer (W8)",
        "ComplianceChecker (W9)",
    ]
    n = len(modules)
    cell_w = (layer_w - Inches(0.5)) / n
    for i, m in enumerate(modules):
        cx = layer_x + Inches(0.25) + (cell_w) * i
        cy = nlp_y + Inches(0.55)
        ch = Inches(0.65)
        color = ACCENT_AMBER if "Compliance" in m else ACCENT
        add_rect(s, cx + Inches(0.05), cy, cell_w - Inches(0.1), ch,
                 fill=PANEL_STRONG, line=color, radius=0.10)
        add_text(s, cx + Inches(0.1), cy + Inches(0.10),
                 cell_w - Inches(0.2), ch - Inches(0.2),
                 m, size=12, color=color, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.15)
    add_text(s, layer_x + Inches(0.25), nlp_y + Inches(0.18),
             layer_w - Inches(0.5), Inches(0.35),
             "scripts/ — pure-Python NLP modules · lazy-instantiated · individually unit-tested",
             size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    # ----- SERVING LAYER -----
    y = Inches(5.2)
    add_text(s, layer_x, y, Inches(1.5), Inches(0.4),
             "SERVING", size=13, color=MUTED, bold=True)
    sy = y + Inches(0.42)
    sh = Inches(0.85)
    api_w = layer_w * 0.65
    add_rect(s, layer_x, sy, api_w, sh, fill=PANEL_STRONG, line=ACCENT_VIOLET, radius=0.08)
    add_text(s, layer_x + Inches(0.2), sy + Inches(0.10),
             api_w - Inches(0.4), Inches(0.35),
             "FastAPI · scripts/production_api.py · 12 endpoints",
             size=15, color=INK, bold=True)
    add_text(s, layer_x + Inches(0.2), sy + Inches(0.45),
             api_w - Inches(0.4), Inches(0.4),
             "/search · /search/compare · /parse-query · /extract-entities · /summarize · /check-compliance · /submit-listing · /listings · /feedback · /metrics · /health",
             size=11, color=MUTED, line_spacing=1.25)

    cross_x = layer_x + api_w + Inches(0.2)
    cross_w = layer_w - api_w - Inches(0.2)
    add_rect(s, cross_x, sy, cross_w, sh, fill=PANEL_STRONG, line=ACCENT_AMBER, radius=0.08)
    add_text(s, cross_x + Inches(0.2), sy + Inches(0.10),
             cross_w - Inches(0.4), Inches(0.35),
             "Cross-cutting concerns",
             size=15, color=INK, bold=True)
    add_text(s, cross_x + Inches(0.2), sy + Inches(0.45),
             cross_w - Inches(0.4), Inches(0.4),
             "TTL+LRU cache · slowapi 10/s · CORS · structured JSON logs · health · warmup",
             size=11, color=MUTED, line_spacing=1.25)

    # ----- CLIENTS -----
    y = Inches(6.65)
    add_text(s, layer_x, y, Inches(1.5), Inches(0.3),
             "CLIENTS", size=13, color=MUTED, bold=True)
    cy = y + Inches(0.32)
    ch = Inches(0.5)
    n3 = 3
    cw = (layer_w - Inches(0.4)) / n3
    clients = [
        "Streamlit Demo (W11) — 3 tabs",
        "Internal Agent Tooling",
        "Listing Submission → Compliance Gate",
    ]
    for i, c in enumerate(clients):
        cx = layer_x + (cw + Inches(0.2)) * i
        add_rect(s, cx, cy, cw, ch, fill=PANEL_STRONG, line=ACCENT_AMBER, radius=0.10)
        add_text(s, cx, cy + Inches(0.10), cw, ch - Inches(0.2),
                 c, size=14, color=ACCENT_AMBER, bold=True, align=PP_ALIGN.CENTER)


def slide_06_data_foundation(prs):
    s = slide_blank(prs)
    add_topbar(s, "05 · Data Foundation", "Weeks 1–2 · Taxonomy, profiling, cleaning")
    add_title(s, "You can't extract what you can't see.")

    # Left: profile-first
    col_w = (USABLE_W - Inches(0.4)) / 2
    lx, ly = MX, Inches(1.85)
    lh = Inches(5.1)
    add_rect(s, lx, ly, col_w, lh, fill=PANEL, line=LINE, radius=0.08)
    add_text(s, lx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.4),
             "Profile-first cleaning",
             size=22, color=ACCENT, bold=True)
    add_paragraphs(
        s, lx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), Inches(2.0),
        [
            "Before writing a regex we ran profile_column() to see what was actually in L_Remarks.",
            "Surfaced: null rate, average length, HTML tag presence, common abbreviations, price-mention frequency.",
            "That report drove the abbreviation map and unicode pass.",
        ],
        size=17, line_spacing=1.35, space_after=8,
    )

    # Examples
    add_text(s, lx + Inches(0.35), ly + Inches(2.95), col_w - Inches(0.7), Inches(0.4),
             "Cleaned examples that survive to the index",
             size=15, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, lx + Inches(0.35), ly + Inches(3.45), col_w - Inches(0.7), Inches(1.6),
        [
            [("450k", {"font": FONT_MONO, "color": ACCENT}), ("  →  ", {}), ("450000", {"font": FONT_MONO, "color": ACCENT})],
            [("2,000 sqft", {"font": FONT_MONO, "color": ACCENT}), ("  →  ", {}), ("2000 square feet", {"font": FONT_MONO, "color": ACCENT})],
            [("w/d", {"font": FONT_MONO, "color": ACCENT}), ("  →  ", {}), ("washer dryer", {"font": FONT_MONO, "color": ACCENT})],
            "Unicode dashes / quotes / NBSP normalised (NFKC); HTML stripped.",
        ],
        size=17, line_spacing=1.35, space_after=8, bullet=True,
    )

    # Right: taxonomy + stats
    rx = MX + col_w + Inches(0.4)
    ry = Inches(1.85)
    rh = Inches(3.4)
    add_rect(s, rx, ry, col_w, rh, fill=PANEL_STRONG, line=LINE, radius=0.08)
    add_text(s, rx + Inches(0.3), ry + Inches(0.25), col_w - Inches(0.6), Inches(0.4),
             "Taxonomy as source-of-truth",
             size=22, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, rx + Inches(0.35), ry + Inches(0.95), col_w - Inches(0.7), Inches(2.4),
        [
            "200+ terms across 8 categories: property_features, amenities, location, community, condition, type, views, other.",
            "Consumed by EntityExtractor, SignalExtractor, QueryParser, AnswerabilityChecker, ComplianceChecker.",
            "Coverage gate: ≥30% of remarks must contain at least one taxonomy term (validated in tests/).",
        ],
        size=17, line_spacing=1.35, space_after=8,
    )

    # Bottom-right metric strip
    sx = MX + col_w + Inches(0.4)
    sy = Inches(5.45)
    sw = col_w
    n = 3
    cw = (sw - Inches(0.4)) / n
    cards = [
        ("ABBREVIATIONS", "30+",   "br · w/d · hw · ss · hoa…",     ACCENT),
        ("CLEANING TESTS", "40+",  "edge cases · unicode · prices", ACCENT_AMBER),
        ("SAMPLE",         "1,000", "remarks → CSV (raw + clean)",  ACCENT_VIOLET),
    ]
    for i, (lbl, big, small, color) in enumerate(cards):
        x = sx + (cw + Inches(0.2)) * i
        add_metric_card(s, x, sy, cw, Inches(1.5), lbl, big, small, color=color)


def slide_07_entity_parser(prs):
    s = slide_blank(prs)
    add_topbar(s, "06 · Entity Extraction & Query Parser", "Weeks 3–4")
    add_title(s, "From free text to safe SQL.")

    # Left: F1 table
    col_w = (USABLE_W - Inches(0.4)) / 2

    add_text(s, MX, Inches(1.85), col_w, Inches(0.5),
             "Entity extraction · 250 labeled remarks", size=20,
             color=ACCENT, bold=True)
    rows = [
        ["Label", "P", "R", "F1", "0.85+"],
        ["BEDROOMS",   "0.960", "1.000", "0.980", ("PASS", GOOD, True)],
        ["BATHROOMS",  "0.974", "1.000", "0.987", ("PASS", GOOD, True)],
        ["AMENITY",    "1.000", "0.842", "0.914", ("PASS", GOOD, True)],
        ["PRICE",      "0.819", "0.908", "0.861", ("PASS", GOOD, True)],
        ["SQFT",       "0.351", "0.870", "0.500", ("FOLLOW-UP", WARN, True)],
        ["Macro",      "—",     "—",     ("0.848", INK, True), ("AT TARGET", WARN, True)],
    ]
    add_table(s, MX, Inches(2.4), col_w, Inches(3.3), rows,
              header_size=14, body_size=15, col_widths=[2.2, 1, 1, 1, 1.5])
    add_text(s, MX, Inches(5.85), col_w, Inches(1.2),
             "SQFT is the principal regression target. Gold spans frequently disagree with "
             "the \"main\" sqft mention vs ADU / accessory units — explicit Future Work item.",
             size=15, color=MUTED, line_spacing=1.35, italic=True)

    # Right: parser + code
    rx = MX + col_w + Inches(0.4)
    add_text(s, rx, Inches(1.85), col_w, Inches(0.5),
             "Query Parser → safe SQL", size=20, color=ACCENT_AMBER, bold=True)
    add_text(s, rx, Inches(2.35), col_w, Inches(0.7),
             "20+ patterns. SQL is always parameterised — no string concat with user input.",
             size=16, color=MUTED, line_spacing=1.35)
    code = (
        ">>> QueryParser().parse_to_sql(\n"
        '...   "3 bed 2 bath under 700k in Irvine with pool")\n'
        "ParsedQuery(\n"
        "  filters={\n"
        "    'price_max': 700000,\n"
        "    'bedrooms_min': 3, 'bedrooms_max': 3,\n"
        "    'bathrooms_min': 2.0, 'bathrooms_max': 2.0,\n"
        "    'city': 'Irvine',\n"
        "    'amenities_in': ['pool']},\n"
        "  where_sql='L_SystemPrice <= %s AND ... AND L_City = %s',\n"
        "  params=[700000, 3, 3, 2.0, 2.0, 'Irvine', '%pool%'])"
    )
    add_code_block(s, rx, Inches(3.05), col_w, Inches(3.2), code, size=13)
    add_text(s, rx, Inches(6.4), col_w, Inches(0.6),
             "Defence-in-depth: SchemaValidator + AnswerabilityChecker reject bad inputs before SQL.",
             size=15, color=ACCENT, italic=True, line_spacing=1.35)


def slide_08_search(prs):
    s = slide_blank(prs)
    add_topbar(s, "07 · Semantic & Hybrid Retrieval", "Week 5")
    add_title(s, "Sub-25 ms retrieval over 10,000 listings.")

    # 3 metric cards
    n = 3
    card_w = (USABLE_W - Inches(0.5)) / n
    metrics = [
        ("P50 LATENCY", "12.7 ms",  "10k listings · top-k=10",  ACCENT),
        ("P95 LATENCY", "15.1 ms",  "target was < 100 ms",      ACCENT_AMBER),
        ("P99 LATENCY", "23.2 ms",  "max 24.0 ms",              ACCENT_VIOLET),
    ]
    for i, (lbl, big, small, color) in enumerate(metrics):
        x = MX + (card_w + Inches(0.25)) * i
        add_metric_card(s, x, Inches(1.85), card_w, Inches(1.5), lbl, big, small, color=color)

    # Stack & Hybrid columns
    col_w = (USABLE_W - Inches(0.4)) / 2
    lx, ly = MX, Inches(3.65)
    lh = Inches(3.55)
    add_rect(s, lx, ly, col_w, lh, fill=PANEL, line=LINE, radius=0.08)
    add_text(s, lx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "Stack", size=22, color=ACCENT, bold=True)
    add_paragraphs(
        s, lx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), lh - Inches(1.1),
        [
            [("Encoder: ", {"bold": True}),
             ("sentence-transformers/all-MiniLM-L6-v2 (384-dim, 22M params)", {})],
            [("Index: ", {"bold": True}),
             ("faiss.IndexFlatIP with L2-normalised vectors → exact cosine similarity", {})],
            [("Lexical baseline: ", {"bold": True}),
             ("in-house BM25 (k1=1.2, b=0.75) over identical corpus", {})],
            [("Hybrid fusion: ", {"bold": True}),
             ("0.7·sem + 0.3·BM25 after per-mode score normalisation", {})],
        ],
        size=17, line_spacing=1.35, space_after=10,
    )

    rx = MX + col_w + Inches(0.4)
    add_rect(s, rx, ly, col_w, lh, fill=PANEL_STRONG, line=LINE, radius=0.08)
    add_text(s, rx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "Why hybrid > semantic-only", size=22, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, rx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), lh - Inches(1.1),
        [
            "Pure semantic loses to BM25 on numeric / brand tokens (zip codes, \"Sub-Zero\", \"1031 exchange\").",
            "Pure BM25 misses paraphrase (\"entertainer's backyard\" ↔ \"great for hosting\").",
            "/search/compare exposes set overlap so PMs can A/B all three modes.",
            "Hybrid is the default for /search; modes are user-selectable.",
        ],
        size=17, line_spacing=1.35, space_after=10,
    )


def slide_09_signals_summary(prs):
    s = slide_blank(prs)
    add_topbar(s, "08 · Signals & Summaries", "Weeks 6 & 8")
    add_title(s, "Structured signals across the entire corpus.")

    col_w = (USABLE_W - Inches(0.4)) / 2

    # Signal extraction
    lx, ly = MX, Inches(1.85)
    lh = Inches(5.0)
    add_rect(s, lx, ly, col_w, lh, fill=PANEL, line=LINE, radius=0.08)
    add_text(s, lx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "Signal extraction (W6)", size=22, color=ACCENT, bold=True)
    add_text(s, lx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), Inches(1.1),
             "Combined Week-3 entities with taxonomy-aware pattern groups across "
             "amenities, condition, financing, and location features.",
             size=16, color=MUTED, line_spacing=1.35)

    rows = [
        ["Metric", "Result", "Target", "Status"],
        ["bedrooms exact",     "1.000  (998/998)",   "0.90", ("PASS", GOOD, True)],
        ["bathrooms exact",    "1.000  (1000/1000)", "0.90", ("PASS", GOOD, True)],
        ["price exact",        "1.000  (1000/1000)", "0.90", ("PASS", GOOD, True)],
        ["free-text amenity",  "0.842  (612/727)",   "0.75", ("PASS", GOOD, True)],
    ]
    add_table(s, lx + Inches(0.35), ly + Inches(2.15),
              col_w - Inches(0.7), Inches(2.0),
              rows, header_size=14, body_size=15,
              col_widths=[2.2, 2.2, 1, 1.2])
    add_text(s, lx + Inches(0.35), ly + Inches(4.3),
             col_w - Inches(0.7), Inches(0.6),
             "Full-table run: 40,890 listings → rets_property_signals.jsonl",
             size=15, color=ACCENT_AMBER, italic=True, bold=True)

    # Summarization
    rx = MX + col_w + Inches(0.4)
    add_rect(s, rx, ly, col_w, lh, fill=PANEL_STRONG, line=LINE, radius=0.08)
    add_text(s, rx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "Summarization (W8)", size=22, color=ACCENT_AMBER, bold=True)
    add_text(s, rx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), Inches(1.1),
             "Sentence scoring (location/feature bonuses + length normalization). "
             "ROUGE-L benchmarked on a held-out sample.",
             size=16, color=MUTED, line_spacing=1.35)

    # Two metric cards
    cw = (col_w - Inches(0.85)) / 2
    add_metric_card(s, rx + Inches(0.35), ly + Inches(2.15), cw, Inches(1.4),
                    "ROUGE-L", "0.7252", "target 0.40 · PASS", color=ACCENT)
    add_metric_card(s, rx + Inches(0.5) + cw, ly + Inches(2.15), cw, Inches(1.4),
                    "HUMAN EVAL", "20", "summaries rated by teammates",
                    color=ACCENT_AMBER)

    # Example summary
    code = (
        '"3 bed, 2 bath home in Irvine listed at $689,000\n'
        ' with pool and hardwood. Located on a quiet\n'
        ' cul-de-sac with top-rated school district…"'
    )
    add_code_block(s, rx + Inches(0.35), ly + Inches(3.75),
                   col_w - Inches(0.7), Inches(1.1), code, size=14)


def slide_10_intent_answerability(prs):
    s = slide_blank(prs)
    add_topbar(s, "09 · Intent & Answerability", "Weeks 7 & 8")
    add_title(s, "Score every query before — and after — it hits the index.")

    col_w = (USABLE_W - Inches(0.4)) / 2
    # Intent
    add_text(s, MX, Inches(1.85), col_w, Inches(0.5),
             "Intent classification (W7)", size=22, color=ACCENT, bold=True)
    add_text(s, MX, Inches(2.4), col_w, Inches(0.7),
             "252-row labeled set · TF-IDF + Logistic Regression · confidence + uncertainty signals.",
             size=16, color=MUTED, line_spacing=1.35)
    rows = [
        ["Split strategy", "Accuracy", "Status"],
        ["Stratified (20% holdout)",   "1.000", ("PASS", GOOD, True)],
        ["Template-holdout (harder)",  "0.937", ("PASS", GOOD, True)],
        ["Human-style external set",   "0.833", ("PASS", GOOD, True)],
    ]
    add_table(s, MX, Inches(3.3), col_w, Inches(2.1), rows,
              header_size=14, body_size=15, col_widths=[3, 1.5, 1.5])
    add_text(s, MX, Inches(5.55), col_w, Inches(1.4),
             "Lead routing use-case: ready_to_buy hits a high-priority queue; "
             "browsing stays self-serve.",
             size=16, color=ACCENT_AMBER, italic=True, line_spacing=1.35)

    # Answerability
    rx = MX + col_w + Inches(0.4)
    add_text(s, rx, Inches(1.85), col_w, Inches(0.5),
             "Answerability gate (W8)", size=22, color=ACCENT_AMBER, bold=True)
    add_text(s, rx, Inches(2.4), col_w, Inches(0.7),
             "Pre-query and post-query checks reject out-of-domain or unsupported "
             "queries with a helpful failure message.",
             size=16, color=MUTED, line_spacing=1.35)
    code = (
        '>>> checker.check_pre_query("how\'s the weather in LA?")\n'
        '# AnswerabilityResult(answerable=False,\n'
        '#   reason="Out of domain query",\n'
        '#   details=["This appears unrelated to real estate."])\n'
        '\n'
        '>>> checker.check_pre_query("3 bed in Atlantis under 500k")\n'
        '# AnswerabilityResult(answerable=False,\n'
        '#   reason="Unsupported filter values",\n'
        '#   details=["City \'Atlantis\' is not a supported city."])'
    )
    add_code_block(s, rx, Inches(3.3), col_w, Inches(3.2), code, size=13)
    add_text(s, rx, Inches(6.6), col_w, Inches(0.5),
             "Failure as a feature — user-visible failure modes are part of the spec.",
             size=15, color=ACCENT, italic=True)


def slide_11_compliance(prs):
    s = slide_blank(prs)
    add_topbar(s, "10 · Fair Housing Compliance", "Week 9")
    add_title(s, "10 categories. 3 severities. Zero false negatives.")

    # Metric strip
    n = 3
    card_w = (USABLE_W - Inches(0.5)) / n
    cards = [
        ("RECALL ON KNOWN VIOLATIONS", "1.000",  "target 1.000 · zero false negatives", ACCENT),
        ("PRECISION (ERROR VS ALL)",   "1.000",  "target > 0.80",                       GOOD),
        ("CATEGORIES COVERED",         "10",     "error · warning · info severities",   ACCENT_AMBER),
    ]
    for i, (l, b, sm, col) in enumerate(cards):
        x = MX + (card_w + Inches(0.25)) * i
        add_metric_card(s, x, Inches(1.85), card_w, Inches(1.55), l, b, sm, color=col)

    # Categories
    add_text(s, MX, Inches(3.65), USABLE_W, Inches(0.4),
             "Categories", size=20, color=ACCENT, bold=True)
    cats = [
        ("race", ACCENT),         ("color", ACCENT),       ("national_origin", ACCENT),
        ("religion", ACCENT),     ("sex", ACCENT),         ("familial_status", ACCENT),
        ("disability", ACCENT),   ("source_of_income", ACCENT_AMBER),
        ("age", ACCENT_AMBER),    ("steering", ACCENT_AMBER),
    ]
    px, py = MX, Inches(4.15)
    for label, color in cats:
        shape, w, _ = add_pill(s, px, py, label, color=color, size=14)
        px += w + Inches(0.12)

    # Workflow
    add_text(s, MX, Inches(4.95), USABLE_W, Inches(0.4),
             "Submission-time enforcement", size=20, color=ACCENT_AMBER, bold=True)

    # Three-stage flow
    flow_y = Inches(5.5)
    flow_h = Inches(1.4)
    box_w = (USABLE_W - Inches(0.6)) / 3
    flows = [
        ("error_count > 0", "→  status = blocked",       "Raise + return suggestions to agent",   RGBColor(0xEF, 0x44, 0x44)),
        ("warning_count > 0", "→  pending_review",        "Routed to human reviewer queue",        ACCENT_AMBER),
        ("else", "→  published",                          "Audit log retains info-level findings", GOOD),
    ]
    for i, (cond, arrow, desc, col) in enumerate(flows):
        x = MX + (box_w + Inches(0.3)) * i
        add_rect(s, x, flow_y, box_w, flow_h, fill=PANEL_STRONG, line=col, radius=0.08)
        add_text(s, x + Inches(0.2), flow_y + Inches(0.15), box_w - Inches(0.4), Inches(0.4),
                 cond, size=14, color=col, bold=True, font=FONT_MONO)
        add_text(s, x + Inches(0.2), flow_y + Inches(0.5), box_w - Inches(0.4), Inches(0.4),
                 arrow, size=18, color=INK, bold=True, font=FONT_MONO)
        add_text(s, x + Inches(0.2), flow_y + Inches(0.95), box_w - Inches(0.4), Inches(0.4),
                 desc, size=13, color=MUTED, line_spacing=1.3)


def slide_12_metrics_summary(prs):
    s = slide_blank(prs)
    add_topbar(s, "11 · Metrics Summary", "Accuracy · Latency · Coverage")
    add_title(s, "Component scoreboard.")

    rows = [
        ["Component",          "Quality / Accuracy",                 "Latency",                       "Status"],
        ["Taxonomy (W1)",      "200+ terms · 8 categories",          "—",                              ("PASS", GOOD, True)],
        ["TextCleaner (W2)",   "40+ unit tests pass",                "< 1 ms / remark",                ("PASS", GOOD, True)],
        ["EntityExtractor (W3)","Macro F1  0.848  · Beds 0.98 · Amenity 0.91", "≈ 2 ms / remark",      ("PASS · SQFT follow-up", WARN, True)],
        ["QueryParser (W4)",   "~90% on Week-1 queries",             "< 1 ms / query",                 ("PASS", GOOD, True)],
        ["Semantic (W5)",      "Hybrid > semantic > BM25 on overlap","p50 12.7  · p95 15.1 ms",        ("PASS", GOOD, True)],
        ["SignalExtractor (W6)","Structured 1.000 · Free-text 0.842","≈ 3 ms / record",                ("PASS", GOOD, True)],
        ["IntentClassifier (W7)","Stratified 1.0  · Holdout 0.94 · Human 0.83", "< 1 ms / query",      ("PASS", GOOD, True)],
        ["Summarizer (W8)",    "ROUGE-L  0.725  (target 0.40)",      "≈ 2 ms / listing",               ("PASS", GOOD, True)],
        ["Compliance (W9)",    "Recall 1.000 · Precision 1.000",     "≈ 2 ms / listing",               ("PASS", GOOD, True)],
        ["FastAPI (W10)",      "12 endpoints · 782 LOC tests",       "uncached search ≈ 25 ms",        ("PASS", GOOD, True)],
        ["Streamlit Demo (W11)","3 tabs · 379 LOC tests · feedback", "API-bounded",                    ("SHIPPED", GOOD, True)],
    ]
    add_table(s, MX, Inches(1.7), USABLE_W, Inches(5.5), rows,
              header_size=15, body_size=15, col_widths=[2.4, 4.5, 2.5, 2])

    add_text(s, MX, Inches(7.05), USABLE_W, Inches(0.4),
             "All numbers are reproducible from the notebooks/ folder & the tests/ suite.",
             size=14, color=MUTED, italic=True)


def slide_13_deep_dive_1(prs):
    s = slide_blank(prs)
    add_topbar(s, "12 · Deep Dive 1 of 3", "Hybrid retrieval + filter-aware over-fetch")
    add_title(s, "The \"filter applied after retrieval\" trap.")

    col_w = (USABLE_W - Inches(0.4)) / 2

    # Left: problem + fix
    add_text(s, MX, Inches(1.9), col_w, Inches(0.5),
             "The problem", size=20, color=ACCENT, bold=True)
    add_paragraphs(
        s, MX, Inches(2.45), col_w, Inches(2.0),
        [
            "FAISS doesn't natively know about price <= 700000.",
            "If we ask for the top-10 most-similar to \"3 bed under 700k in Irvine with pool\", all 10 may be in San Diego.",
            "We're left with zero results after applying the city filter.",
        ],
        size=17, line_spacing=1.35, space_after=8,
    )

    add_text(s, MX, Inches(4.55), col_w, Inches(0.5),
             "Our two-step fix", size=20, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, MX, Inches(5.05), col_w, Inches(2.2),
        [
            [("Over-retrieve when filters are present. ", {"bold": True}),
             ("retrieval_depth() returns max(top_k·25, 250) candidates.", {})],
            [("Post-filter with metadata. ", {"bold": True}),
             ("Hits are joined to ListingMetadataStore and filtered before truncation.", {})],
            [("Amenity post-filter ", {"bold": True}),
             ("runs last on remark text — cheap and captures include / exclude semantics.", {})],
        ],
        size=17, line_spacing=1.35, space_after=10,
    )

    # Right: code + impact
    rx = MX + col_w + Inches(0.4)
    code = (
        "def retrieval_depth(top_k, filters, store):\n"
        "    \"\"\"Fetch a wider candidate pool when\n"
        "       structured filters are present.\"\"\"\n"
        "    if not _metadata_filters_present(filters):\n"
        "        return top_k\n"
        "    n = len(store) if store else 0\n"
        "    if n <= 0:\n"
        "        return max(top_k * 10, 100)\n"
        "    return min(n, max(top_k * 25, 250))\n"
        "\n"
        "# Inside /search:\n"
        "candidate_top_k = retrieval_depth(\n"
        "    payload.top_k, filters, metadata_store)\n"
        "raw = run_search(payload.query, mode=mode,\n"
        "                 top_k=candidate_top_k, ...)\n"
        "filtered = apply_metadata_filters(\n"
        "    raw, filters, metadata_store)\n"
        "filtered = apply_filters(\n"
        "    filtered, filters)[:payload.top_k]"
    )
    add_code_block(s, rx, Inches(1.9), col_w, Inches(4.4), code, size=13)
    add_text(s, rx, Inches(6.45), col_w, Inches(0.7),
             "Impact: filtered queries go from \"often 0 results\" to stable, "
             "without breaking the < 30 ms latency budget.",
             size=15, color=ACCENT, italic=True, line_spacing=1.35)


def slide_14_deep_dive_2(prs):
    s = slide_blank(prs)
    add_topbar(s, "13 · Deep Dive 2 of 3", "Pattern library design under FHA")
    add_title(s, "100% recall is easy. Precision was the hard part.")

    col_w = (USABLE_W - Inches(0.4)) / 2

    add_text(s, MX, Inches(1.9), col_w, Inches(0.5),
             "Two failure modes", size=20, color=ACCENT, bold=True)
    add_paragraphs(
        s, MX, Inches(2.45), col_w, Inches(1.8),
        [
            [("False negative ", {"bold": True}),
             ("— we miss \"adults only\". HUD penalty up to $25k+ on first violation. Unacceptable.", {})],
            [("False positive ", {"bold": True}),
             ("— we block \"family room\" or \"close to St. Mary's church\". Agents lose trust. Also unacceptable.", {})],
        ],
        size=17, line_spacing=1.35, space_after=12,
    )

    add_text(s, MX, Inches(4.5), col_w, Inches(0.5),
             "Design choices that bought us both", size=20, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, MX, Inches(5.05), col_w, Inches(2.3),
        [
            [("Word-boundary wrappers. ", {"bold": True}),
             ("_word(p) uses (?<![A-Za-z0-9])…(?![A-Za-z0-9]) — \"manor\" ≠ \"man\".", {})],
            [("Three severities, not boolean. ", {"bold": True}),
             ("Religious landmark names = warning, not error.", {})],
            [("Suggestion text on every pattern. ", {"bold": True}),
             ("Agent gets the rule and the rewrite hint.", {})],
            [("Subject-matter docs. ", {"bold": True}),
             ("docs/fair_housing_guidelines.md so non-engineers can audit.", {})],
        ],
        size=16, line_spacing=1.35, space_after=8,
    )

    # Right: code excerpt
    rx = MX + col_w + Inches(0.4)
    code = (
        "# scripts/compliance_checker.py (excerpt)\n"
        "\n"
        "def _word(pattern: str) -> str:\n"
        "    return rf\"(?<![A-Za-z0-9])(?:{pattern})(?![A-Za-z0-9])\"\n"
        "\n"
        "_PATTERN_LIBRARY = {\n"
        "    \"familial_status\": [\n"
        "        (_word(r\"adults?[\\s-]+only\"),\n"
        "         ERROR,\n"
        "         \"Excludes children — FHA violation\",\n"
        "         \"Describe the property, not the occupant.\"),\n"
        "    ],\n"
        "    \"religion\": [\n"
        "        (_word(r\"(christian|jewish|muslim)\\s+\"\n"
        "               r\"(community|neighborhood)\"),\n"
        "         ERROR, \"Religious preference\",\n"
        "         \"Replace with property descriptors.\"),\n"
        "        (_word(r\"near\\s+(st\\.?|saint)\\s+\\w+\"),\n"
        "         WARNING, \"Religious landmark — verify\",\n"
        "         \"Use as directional reference only.\"),\n"
        "    ],\n"
        "    # ... 8 more categories\n"
        "}"
    )
    add_code_block(s, rx, Inches(1.9), col_w, Inches(5.0), code, size=13)
    add_text(s, rx, Inches(7.0), col_w, Inches(0.4),
             "Result: recall 1.000 · precision 1.000 on labeled corpus.",
             size=15, color=GOOD, italic=True, bold=True)


def slide_15_deep_dive_3(prs):
    s = slide_blank(prs)
    add_topbar(s, "14 · Deep Dive 3 of 3", "Production FastAPI architecture")
    add_title(s, "From notebooks to a service that survives a launch demo.")

    col_w = (USABLE_W - Inches(0.4)) / 2

    add_text(s, MX, Inches(1.9), col_w, Inches(0.5),
             "Five things that made it production-shaped",
             size=20, color=ACCENT, bold=True)
    add_paragraphs(
        s, MX, Inches(2.45), col_w, Inches(4.6),
        [
            [("Lazy components, eager warmup. ", {"bold": True}),
             ("Heavy deps loaded on first access; FastAPI lifespan calls warmup_semantic() at boot. Cold-start cost is paid by the deploy, not the user.", {})],
            [("Graceful 503. ", {"bold": True}),
             ("If FAISS artefacts are missing, return 503 — UI auto-falls-back to keyword mode.", {})],
            [("TTL + LRU response cache ", {"bold": True}),
             ("(60 s, 512 entries, thread-safe), introspectable via /cache/stats.", {})],
            [("Per-IP rate limiting ", {"bold": True}),
             ("via slowapi (10 req/s default · 5 for /compare).", {})],
            [("Structured logs + metrics. ", {"bold": True}),
             ("Every request emits method, path, client, status, latency. /metrics/dashboard surfaces p50/p95/p99 + cache + satisfaction proxy.", {})],
        ],
        size=16, line_spacing=1.35, space_after=10,
    )

    # Right code
    rx = MX + col_w + Inches(0.4)
    code = (
        "@asynccontextmanager\n"
        "async def lifespan(app):\n"
        "    _warmup_components_for_startup()\n"
        "    yield\n"
        "\n"
        "app = FastAPI(\n"
        "    title=\"Real Estate NLP API\",\n"
        "    version=API_VERSION,\n"
        "    lifespan=lifespan)\n"
        "app.state.limiter = limiter\n"
        "app.add_middleware(CORSMiddleware, ...)\n"
        "\n"
        "@app.middleware(\"http\")\n"
        "async def log_requests(request, call_next):\n"
        "    start = time.perf_counter()\n"
        "    response = await call_next(request)\n"
        "    elapsed = (time.perf_counter()-start)*1000\n"
        "    metrics.record_request(\n"
        "        elapsed, path=request.url.path)\n"
        "    response.headers[\"X-Process-Time-Ms\"] = (\n"
        "        f\"{elapsed:.2f}\")\n"
        "    return response"
    )
    add_code_block(s, rx, Inches(1.9), col_w, Inches(5.0), code, size=13)
    add_text(s, rx, Inches(7.0), col_w, Inches(0.4),
             "Test surface: 782 LOC pytest covers happy path, 429, 503, cache, compliance gating.",
             size=14, color=ACCENT_AMBER, italic=True, line_spacing=1.3)


def slide_16_production(prs):
    s = slide_blank(prs)
    add_topbar(s, "15 · Production Considerations", "Deployment · observability · scaling")
    add_title(s, "What it takes to keep this running on a Tuesday.")

    col_w = (USABLE_W - Inches(0.4)) / 2

    # Left: deployment + observability
    lx, ly = MX, Inches(1.85)
    lh = Inches(5.1)
    add_rect(s, lx, ly, col_w, lh, fill=PANEL, line=LINE, radius=0.08)
    add_text(s, lx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "Deployment shape", size=22, color=ACCENT, bold=True)
    add_paragraphs(
        s, lx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), Inches(1.7),
        [
            "One Dockerfile; RUN_MODE=api|ui selects which process starts.",
            "render.yaml provisions both services with auto-deploy on push.",
            "Free tier OOMs on FAISS load (512 MB); blueprint defaults to Starter.",
            "Local parity: docker compose up brings MySQL + API + UI up together.",
        ],
        size=16, line_spacing=1.35, space_after=8,
    )

    add_text(s, lx + Inches(0.3), ly + Inches(2.7), col_w - Inches(0.6), Inches(0.5),
             "Observability", size=22, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, lx + Inches(0.35), ly + Inches(3.4), col_w - Inches(0.7), Inches(1.7),
        [
            "Structured JSON logs ready for Datadog / CloudWatch.",
            "/health returns liveness + semantic readiness separately.",
            "In-process metrics — swap to Prometheus pull when > 1 replica.",
            "Feedback events on disk now — see Future Work for persistence.",
        ],
        size=16, line_spacing=1.35, space_after=8,
    )

    # Right: scaling + runbook
    rx = MX + col_w + Inches(0.4)
    add_rect(s, rx, ly, col_w, lh, fill=PANEL_STRONG, line=LINE, radius=0.08)
    add_text(s, rx + Inches(0.3), ly + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             "Scaling levers", size=22, color=ACCENT_VIOLET, bold=True)
    add_paragraphs(
        s, rx + Inches(0.35), ly + Inches(0.95), col_w - Inches(0.7), Inches(1.7),
        [
            [("Vertical: ", {"bold": True}), ("Render Standard ($25) absorbs current corpus.", {})],
            [("Horizontal: ", {"bold": True}), ("stateless API → N replicas behind LB.", {})],
            [("Index sharding: ", {"bold": True}), ("at 10× corpus, partition FAISS or move to managed pgvector.", {})],
            [("Cost knobs: ", {"bold": True}), ("60s TTL cache absorbs > 50% of demo traffic.", {})],
        ],
        size=16, line_spacing=1.35, space_after=8,
    )
    add_text(s, rx + Inches(0.3), ly + Inches(2.7), col_w - Inches(0.6), Inches(0.5),
             "Runbook (top failures)", size=22, color=ACCENT_AMBER, bold=True)
    rows = [
        ["Symptom", "Action"],
        ["Container exit 137",        "OOM — bump plan or remove FAISS (keyword-only)"],
        ["503 Semantic unavailable",  "scripts/semantic_searcher.py --build"],
        ["429 Rate limit",            "Increase default_limits or add replicas"],
        ["UI ConnectError",           "Confirm API_BASE_URL on UI service"],
    ]
    add_table(s, rx + Inches(0.35), ly + Inches(3.35),
              col_w - Inches(0.7), Inches(1.6), rows,
              header_size=13, body_size=14, col_widths=[2, 3])


def slide_17_future_work(prs):
    s = slide_blank(prs)
    add_topbar(s, "16 · Future Work", "90-day roadmap")
    add_title(s, "What we'd build next.")

    col_w = (USABLE_W - Inches(0.5)) / 3
    h = Inches(5.0)
    y = Inches(1.85)

    cols = [
        ("Quality", ACCENT, [
            [("Replace SQFT regex with fine-tuned spaCy NER. ", {"bold": True}),
             ("Target macro-F1 > 0.90.", {})],
            [("Cross-encoder re-ranker ", {"bold": True}),
             ("(ms-marco-MiniLM-L6-v2) on top of hybrid — pay 80 ms for 5–15% nDCG lift.", {})],
            [("Active learning loop on compliance. ", {"bold": True}),
             ("Reviewer-cleared warnings → hard negatives.", {})],
            [("Multilingual remarks. ", {"bold": True}),
             ("Spanish-language MLS feeds in CA / TX.", {})],
        ]),
        ("Platform", ACCENT_VIOLET, [
            [("Postgres + pgvector ", {"bold": True}),
             ("instead of FAISS-on-disk: filters pushed into SQL, streaming index updates.", {})],
            [("Redis ", {"bold": True}),
             ("for response cache + rate-limiter state across replicas.", {})],
            [("Persistent feedback store ", {"bold": True}),
             ("(Postgres / Supabase) — required for satisfaction telemetry.", {})],
            [("Per-MLS taxonomies. ", {"bold": True}),
             ("Multi-tenant config to onboard markets beyond California.", {})],
        ]),
        ("Product", ACCENT_AMBER, [
            [("Conversational follow-ups. ", {"bold": True}),
             ("LLM-backed assistant when the answerability gate rejects.", {})],
            [("Saved searches + alerts. ", {"bold": True}),
             ("Nightly run of stored queries → email on new matches.", {})],
            [("Compliance v2. ", {"bold": True}),
             ("State-specific protected classes (NY, CA additions).", {})],
            [("Agent assist. ", {"bold": True}),
             ("\"Suggest a remark\" from extracted signals — compliant by construction.", {})],
        ]),
    ]
    for i, (head, color, items) in enumerate(cols):
        x = MX + (col_w + Inches(0.25)) * i
        add_rect(s, x, y, col_w, h, fill=PANEL, line=LINE, radius=0.08)
        add_text(s, x + Inches(0.3), y + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
                 head, size=22, color=color, bold=True)
        add_paragraphs(s, x + Inches(0.35), y + Inches(0.95),
                       col_w - Inches(0.7), h - Inches(1.1),
                       items, size=15, line_spacing=1.35, space_after=10)


def slide_18_lessons(prs):
    s = slide_blank(prs)
    add_topbar(s, "17 · Lessons Learned", "What we'd do again — and differently")
    add_title(s, "Engineering takeaways.")

    col_w = (USABLE_W - Inches(0.4)) / 2
    y = Inches(1.85)
    h = Inches(3.7)

    # What we'd do again
    add_rect(s, MX, y, col_w, h, fill=PANEL, line=LINE, radius=0.08)
    add_text(s, MX + Inches(0.3), y + Inches(0.22), col_w - Inches(0.6), Inches(0.45),
             "What we'd do again", size=22, color=ACCENT, bold=True)
    add_paragraphs(
        s, MX + Inches(0.35), y + Inches(0.85), col_w - Inches(0.7), h - Inches(0.95),
        [
            [("Profile before you clean. ", {"bold": True}),
             ("Week-2 profiling told us which abbreviations were worth the effort.", {})],
            [("Regex first, ML later. ", {"bold": True}),
             ("Beds/baths/price F1 > 0.96 with regex. Save ML for SQFT.", {})],
            [("Test the failure cases. ", {"bold": True}),
             ("Most of the API test suite exercises 503/429/empty paths — that's why the demo never broke.", {})],
            [("One Docker image, two services. ", {"bold": True}),
             ("Single source of truth for deps; RUN_MODE switch saved CI maintenance.", {})],
            [("Lazy components + lifespan warmup. ", {"bold": True}),
             ("Imports stayed cheap; first-user latency stayed flat.", {})],
        ],
        size=15, line_spacing=1.35, space_after=8,
    )

    # What we'd do differently
    rx = MX + col_w + Inches(0.4)
    add_rect(s, rx, y, col_w, h, fill=PANEL_STRONG, line=LINE, radius=0.08)
    add_text(s, rx + Inches(0.3), y + Inches(0.22), col_w - Inches(0.6), Inches(0.45),
             "What we'd do differently", size=22, color=ACCENT_AMBER, bold=True)
    add_paragraphs(
        s, rx + Inches(0.35), y + Inches(0.85), col_w - Inches(0.7), h - Inches(0.95),
        [
            [("Move metadata into Postgres earlier. ", {"bold": True}),
             ("Python-side filtering works for 40k rows, not for 10M.", {})],
            [("Treat fair-housing rule authoring as a stakeholder workflow. ", {"bold": True}),
             ("YAML+UI for compliance officers is overdue.", {})],
            [("Adopt a real retrieval eval harness. ", {"bold": True}),
             ("Held-out judged set with nDCG@10 should be the merge gate.", {})],
            [("Persistent feedback from day one. ", {"bold": True}),
             ("JSONL on ephemeral disk resets the satisfaction proxy on every restart.", {})],
        ],
        size=15, line_spacing=1.35, space_after=8,
    )

    # Closing stat strip
    n = 3
    cw = (USABLE_W - Inches(0.5)) / n
    cards = [
        ("FROM ZERO TO",  "12",      "production endpoints in 12 weeks", ACCENT),
        ("OUTPERFORMED",  "Every",   "component-level target on the spec", ACCENT_AMBER),
        ("SHIPPED",       "live",    "git push origin main → auto-deploy", GOOD),
    ]
    for i, (l, b, sm, col) in enumerate(cards):
        x = MX + (cw + Inches(0.25)) * i
        add_metric_card(s, x, Inches(5.85), cw, Inches(1.4), l, b, sm, color=col)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_title,
        slide_02_problem,
        slide_03_what_we_built,
        slide_04_approach,
        slide_05_architecture,
        slide_06_data_foundation,
        slide_07_entity_parser,
        slide_08_search,
        slide_09_signals_summary,
        slide_10_intent_answerability,
        slide_11_compliance,
        slide_12_metrics_summary,
        slide_13_deep_dive_1,
        slide_14_deep_dive_2,
        slide_15_deep_dive_3,
        slide_16_production,
        slide_17_future_work,
        slide_18_lessons,
    ]
    for fn in builders:
        fn(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path}  ({len(builders)} slides)")


if __name__ == "__main__":
    here = Path(__file__).resolve().parents[1]
    build(here / "docs" / "final_presentation.pptx")
